#!/usr/bin/env python3
"""Compose a 'notes pages' PDF (slide image on top, speaker notes below) from the
current .pptx and pre-rendered slide PNGs in _notes_tmp/.  Reflects whatever is saved
in the .pptx right now.  Run after rendering slides to _notes_tmp/."""
import os, glob
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

ROOT = r"c:/Users/Administrator/smart-elevator-twin"
PPTX = os.path.join(ROOT, "Agentic_Elevator_Digital_Twin_Defense.pptx")
TMP  = os.path.join(ROOT, "_notes_tmp")
OUT  = os.path.join(ROOT, "Agentic_Elevator_Digital_Twin_Defense_NOTES.pdf")

W, H = 1240, 1754                      # ~A4 portrait @150dpi
M = 64                                 # margin
INK = (15, 23, 42); BODY = (31, 41, 55); SUB = (100, 116, 139); CYAN = (14, 136, 153)
RULE = (226, 232, 240); GREYTXT = (148, 163, 184)

def font(path, size):
    return ImageFont.truetype(path, size)
FB = r"C:/Windows/Fonts/segoeuib.ttf"  # bold
FR = r"C:/Windows/Fonts/segoeui.ttf"   # regular
FI = r"C:/Windows/Fonts/segoeuii.ttf"  # italic

prs = Presentation(PPTX)
slides = list(prs.slides)

def slide_title(s):
    """Best-effort short title: kicker + title, skipping the footer line."""
    texts = []
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t and "Agentic AI-Driven Elevator Digital Twin" not in t and not t.isdigit():
                texts.append(t.replace("\n", " "))
    if not texts:
        return ""
    # kicker is usually short upper-case; title is the next line
    head = texts[0]
    if len(texts) > 1 and len(texts[0]) < 40 and texts[0].upper() == texts[0]:
        head = f"{texts[0]} — {texts[1]}"
    return head[:90]

def wrap(draw, text, fnt, maxw):
    out = []
    for para in text.split("\n"):
        para = para.strip()
        if not para:
            out.append("")
            continue
        words = para.split(" ")
        line = ""
        for w in words:
            trial = (line + " " + w).strip()
            if draw.textlength(trial, font=fnt) <= maxw:
                line = trial
            else:
                if line:
                    out.append(line)
                line = w
        if line:
            out.append(line)
    return out

pages = []
for i, s in enumerate(slides, 1):
    png = os.path.join(TMP, f"s{i:02d}.png")
    page = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(page)
    # top accent + header
    d.rectangle([0, 0, W, 8], fill=CYAN)
    d.text((M, 34), f"Slide {i} of {len(slides)}", font=font(FB, 20), fill=CYAN)
    title = slide_title(s)
    if title:
        for j, ln in enumerate(wrap(d, title, font(FB, 24), W - 2*M)[:2]):
            d.text((M, 60 + j*30), ln, font=font(FB, 24), fill=INK)
    # slide image
    iy = 128
    if os.path.exists(png):
        im = Image.open(png).convert("RGB")
        iw = W - 2*M
        ih = int(iw * im.height / im.width)
        im = im.resize((iw, ih))
        page.paste(im, (M, iy))
        d.rectangle([M, iy, M+iw, iy+ih], outline=RULE, width=2)
        notes_top = iy + ih + 34
    else:
        notes_top = iy
    # notes label
    d.line([M, notes_top, W-M, notes_top], fill=RULE, width=2)
    d.text((M, notes_top + 14), "SPEAKER NOTES", font=font(FB, 20), fill=SUB)
    body_top = notes_top + 50
    raw = ""
    if s.has_notes_slide:
        raw = s.notes_slide.notes_text_frame.text.strip()
    avail = H - M - body_top
    if not raw:
        d.text((M, body_top), "(no speaker notes)", font=font(FI, 24), fill=GREYTXT)
    else:
        # fit font size so wrapped notes fit the available height
        for size in (26, 24, 22, 20, 18, 16, 14):
            fnt = font(FR, size); lh = int(size * 1.32)
            lines = wrap(d, raw, fnt, W - 2*M)
            if len(lines) * lh <= avail or size == 14:
                break
        y = body_top
        for ln in lines:
            if y + lh > H - M:
                break
            d.text((M, y), ln, font=fnt, fill=BODY)
            y += lh
    pages.append(page)

pages[0].save(OUT, "PDF", save_all=True, append_images=pages[1:], resolution=150.0)
print("Saved:", OUT, "| pages:", len(pages))
