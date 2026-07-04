#!/usr/bin/env python3
"""
Defence deck builder v2 — "Agentic AI-Driven Digital Twin for Smart and Secure
Elevator Management System". Two presenters, light theme, simple+technical English,
"we" voice, synoptic part-by-part deep-dive, dedicated ML slide, integration vision.
Reuses the authors' own assets (brand logos, 3D render, synoptic, tech icons) extracted
to _pptx_media/. Run:  python scripts/build_defense_deck.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
ROOT = r"c:/Users/Administrator/smart-elevator-twin"
FIG  = r"c:/Users/Administrator/OneDrive/Desktop/abderrahmane/master thesis dissertation/master-thesis1/figures"
MEDIA = os.path.join(ROOT, "_pptx_media")
OUT  = os.path.join(ROOT, "Agentic_Elevator_Digital_Twin_Defense.pptx")
def F(p): return os.path.join(FIG, p)
def M(p): return os.path.join(MEDIA, p)

# Authors' own assets (extracted from their pptx)
RENDER   = M("image7.jpeg")     # 3D elevator render
UNIV     = M("image1.png")      # university logo
SYNOPTIC = M("image39.png")     # polished labelled synoptic
LOGO = {  # title-slide brand logos (x,y,w,h replicated from their layout)
    "next":  (M("image6.png"), 8.85, 0.56, 0.88, 0.88),
    "mosq":  (M("image2.png"), 8.81, 1.58, 0.95, 0.93),
    "ditto": (M("image3.png"), 8.83, 2.70, 0.92, 1.14),
    "docker":(M("image5.png"), 8.72, 4.07, 1.13, 0.82),
    "n8n":   (M("image4.png"), 8.39, 5.15, 1.81, 0.46),
    "ollama":(M("image8.png"), 8.53, 5.76, 1.49, 0.57),
    "pg":    (M("image9.png"), 8.54, 6.51, 1.52, 0.78),
}
# tech-slide card icons, in card order
TECH_ICON = [M(f"image{n}.png") for n in (17,19,21,23,25,27,29,31,33,35)]
# brand logos reused as small accents on deep-dive slides
BRAND = {"mosq":M("image2.png"),"ditto":M("image3.png"),"n8n":M("image4.png"),
         "docker":M("image5.png"),"next":M("image6.png"),"ollama":M("image8.png"),"pg":M("image9.png")}

# --------------------------------------------------------------------------- #
BG=RGBColor(0xF8,0xFA,0xFC); CARD=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x0F,0x17,0x2A)
SUBTLE=RGBColor(0x47,0x55,0x69); CYAN=RGBColor(0x12,0xC6,0xD8); CYAN_DK=RGBColor(0x0E,0x88,0x99)
STEEL=RGBColor(0x64,0x72,0x7A); RULE=RGBColor(0xE2,0xE8,0xF0); RED=RGBColor(0xEF,0x4D,0x43)
RED_BG=RGBColor(0xFE,0xE2,0xE0); GREEN=RGBColor(0x16,0xA3,0x4A); GREEN_BG=RGBColor(0xDC,0xFC,0xE7)
AMBER=RGBColor(0xB4,0x53,0x09); AMBER_BG=RGBColor(0xFE,0xF3,0xC7); GREY=RGBColor(0x64,0x74,0x8B)
GREY_BG=RGBColor(0xF1,0xF5,0xF9); CYAN_BG=RGBColor(0xE0,0xF7,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)

import glob as _glob
_has_inter = bool(_glob.glob(r"C:/Windows/Fonts/Inter*.ttf"))
FONT = "Inter" if _has_inter else "Segoe UI"
MONO = "Consolas"
A1, A2 = "Abderrahmane Zaouidi", "Mohamed Hachem Bengherbia"
AUTHORS = f"{A1} – {A2}"
SUP = "Pr. Mounir Bouhedda"
SHORT = "Agentic AI-Driven Elevator Digital Twin"

prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH = prs.slide_width, prs.slide_height; BLANK = prs.slide_layouts[6]

# --------------------------------------------------------------------------- helpers
def _noshadow(sh):
    try: sh.shadow.inherit=False
    except Exception: pass

def slide():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=BG; r.line.fill.background(); _noshadow(r)
    sp=r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
    return s

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,radius=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    _noshadow(sp)
    if radius is not None and shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    return sp

def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=4,line_spacing=1.05,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0)
        try: p.line_spacing=line_spacing
        except Exception: pass
        for (txt,size,color,bold,font,italic) in para:
            r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
            r.font.italic=italic; r.font.name=font; r.font.color.rgb=color
    return tb

def P(t,sz,c=INK,b=False,f=None,it=False): return (t,sz,c,b,f or FONT,it)

def chip(s,x,y,label,fill,fg,size=10.5,h=0.30,w=None):
    w=w or (0.115*len(label)+0.26)
    sp=rect(s,x,y,w,h,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    tf=sp.text_frame; tf.word_wrap=False
    tf.margin_left=Inches(0.07); tf.margin_right=Inches(0.07); tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=label; r.font.size=Pt(size); r.font.bold=True; r.font.name=FONT; r.font.color.rgb=fg
    return x+w

def footer(s,page):
    text(s,0.6,7.07,9.5,0.3,[[P(SHORT+"   •   "+A1+" - "+A2,9,STEEL)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,11.9,7.07,0.83,0.3,[[P(str(page),9,STEEL,True)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)

def title_block(s,kicker,ttl,page):
    rect(s,0,0,13.333,0.06,fill=CYAN)
    text(s,0.6,0.42,12,0.28,[[P(kicker.upper(),12,CYAN_DK,True)]])
    text(s,0.6,0.70,12.1,0.8,[[P(ttl,26,INK,True)]])
    rect(s,0.6,1.52,1.7,0.045,fill=CYAN)
    footer(s,page)

def notes(s,t): s.notes_slide.notes_text_frame.text=t.strip()

def place(s,path,x,y,w,h,border=False):
    pic=s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    if border: pic.line.color.rgb=RULE; pic.line.width=Pt(1.0)
    return pic

def picture(s,path,x,y,w,h,caption=None,border=True,align="center",valign="middle"):
    iw,ih=Image.open(path).size; ratio=min(w/(iw/96.0),h/(ih/96.0))
    dw=(iw/96.0)*ratio; dh=(ih/96.0)*ratio
    px=x+(w-dw)/2 if align=="center" else (x if align=="left" else x+(w-dw))
    py=y+(h-dh)/2 if valign=="middle" else (y if valign=="top" else y+(h-dh))
    pic=s.shapes.add_picture(path,Inches(px),Inches(py),Inches(dw),Inches(dh))
    if border: pic.line.color.rgb=RULE; pic.line.width=Pt(1.0)
    if caption: text(s,x,y+h+0.02,w,0.4,[[P(caption,10,SUBTLE,it=True)]],align=PP_ALIGN.CENTER)
    return pic

def arrow(s,x,y,w,h,fill=STEEL,shape=MSO_SHAPE.RIGHT_ARROW):
    return rect(s,x,y,w,h,fill=fill,shape=shape)

# pipeline strip for the deep-dive section
STAGES=["Field","Messaging","Twin","Agentic","Dispatch","Safety","Data","Screen"]
def pipeline(s,active,y=1.66):
    n=len(STAGES); total=12.13; gap=0.08; w=(total-(n-1)*gap)/n; x=0.6
    for i,name in enumerate(STAGES):
        bx=x+i*(w+gap); act=(i==active)
        sp=rect(s,bx,y,w,0.46,fill=(CYAN if act else GREY_BG),
                shape=MSO_SHAPE.PENTAGON if i<n-1 else MSO_SHAPE.CHEVRON, radius=None)
        try: sp.adjustments[0]=0.4
        except Exception: pass
        tf=sp.text_frame; tf.word_wrap=False; tf.margin_top=0; tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=name; r.font.size=Pt(9.5); r.font.bold=act
        r.font.name=FONT; r.font.color.rgb=(WHITE if act else STEEL)

def wwhw(s,what,which,how,why,x=0.6,y=2.55,w=12.13,ch=1.85):
    cw=(w-0.25)/2
    cells=[("WHAT we did",what,CYAN_DK,CYAN_BG),("WHICH tools",which,STEEL,GREY_BG),
           ("HOW it works",how,GREEN,GREEN_BG),("WHY this choice",why,AMBER,AMBER_BG)]
    for i,(lab,body,col,bg) in enumerate(cells):
        bx=x+(i%2)*(cw+0.25); by=y+(i//2)*(ch+0.2)
        rect(s,bx,by,cw,ch,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
        rect(s,bx,by,0.10,ch,fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
        text(s,bx+0.28,by+0.13,cw-0.45,0.3,[[P(lab,11.5,col,True)]])
        text(s,bx+0.28,by+0.48,cw-0.45,ch-0.6,[[P(body,12.5,INK)]])

# =========================================================================== #
# 1 — TITLE
# =========================================================================== #
s=slide()
rect(s,0,0,0.22,7.5,fill=CYAN)
place(s,RENDER,10.2,-0.01,3.16,7.52)
place(s,UNIV,0.7,0.55,1.25,1.2)
for k,(path,x,y,w,h) in LOGO.items(): place(s,path,x,y,w,h)
text(s,0.7,2.0,7.6,0.4,[[P("MASTER'S THESIS DEFENCE",13,CYAN_DK,True)]])
text(s,0.7,2.42,7.9,1.9,[[P("Agentic AI-Driven Digital Twin",31,INK,True)],
                         [P("for Smart and Secure Elevator",26,INK,True)],
                         [P("Management System",26,INK,True)]],line_spacing=1.04)
rect(s,0.72,4.5,2.0,0.05,fill=CYAN)
text(s,0.7,4.7,7.8,1.0,[[P("A software platform for smart, secure elevator supervision — with a",13.5,SUBTLE)],
                        [P("reduced-scale ESP32-S3 elevator as our test bench.",13.5,SUBTLE)]])
text(s,0.7,5.8,7.8,1.0,[[P("Presented by:  ",12.5,STEEL,True),P(AUTHORS,12.5,INK,True)],
                        [P("Supervisor:  ",12.5,STEEL,True),P(SUP,12.5,INK)]])
text(s,0.7,7.0,3,0.3,[[P("2026",12,STEEL,True)]])
notes(s,"""
Good morning. We are Abderrahmane Zaouidi and Mohamed Hachem Bengherbia, and we present our
thesis: an Agentic AI-Driven Digital Twin for Smart and Secure Elevator Management. One
framing to keep in mind from the start: this is first a software platform. The small
four-floor elevator we built is our test bench - it lets us prove the software safely and
cheaply. The same software is designed to drive a real elevator. Over the next twenty
minutes we will show the problem, our solution, the platform part by part, the
machine-learning model we are building, and what we honestly validated. [~35s] -> Here is
our plan.
""")

# =========================================================================== #
# 2 — PLAN / AGENDA
# =========================================================================== #
s=slide(); title_block(s,"Plan","What we will cover",2)
items=[("1","Context & problem — elevators today"),
       ("2","Our solution — and why it is better"),
       ("3","The technology, in plain words"),
       ("4","Inside the platform — architecture, part by part"),
       ("5","The machine-learning model we are building"),
       ("6","Validation & evidence"),
       ("7","Limits, future work & the bigger picture"),
       ("8","Conclusion")]
y=1.95
for i,(n,t) in enumerate(items):
    col=0 if i<4 else 1; row=i%4
    bx=0.6+col*6.17; by=1.95+row*1.18
    rect(s,bx,by,5.95,1.0,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    o=rect(s,bx+0.2,by+0.22,0.56,0.56,fill=CYAN_BG,shape=MSO_SHAPE.OVAL)
    tf=o.text_frame; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=n; r.font.bold=True; r.font.size=Pt(17); r.font.name=FONT; r.font.color.rgb=CYAN_DK
    text(s,bx+0.95,by+0.16,5.0,0.7,[[P(t,14,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
notes(s,"""
Here is our plan. We start with how elevators are managed today and what is wrong with it.
Then we present our solution and why it beats a simple old elevator. We explain the
technology in plain words, then we open the platform and walk it part by part. We give a
dedicated slide to the machine-learning model we are building. Then our validation and
honest evidence, the limits and the bigger picture, and a short conclusion. [~25s] -> Let
us start with the problem.
""")

# =========================================================================== #
# 3 — CONTEXT, PROBLEM & STATE OF THE ART
# =========================================================================== #
s=slide(); title_block(s,"Context & Problem","Elevators today react late, cost a lot, and stay isolated",3)
text(s,0.6,1.62,12.1,0.4,[[P("State of the art: ",12,STEEL,True),
    P("today's elevators use fixed, isolated controllers — no live digital copy, no audited remote commands, security added late.",12,SUBTLE)]])
pains=[("Reactive maintenance","Service is on a fixed schedule; faults are seen only after they happen."),
       ("Long waiting time","Fixed, clock-based dispatch (cabin strategy) wastes time at peak hours."),
       ("High electricity bills","Elevators run hard with no energy strategy."),
       ("High maintenance costs","Unplanned breakdowns and emergency call-outs are expensive."),
       ("Weak, siloed security","Access control and audit are limited and not joined up."),
       ("Fragmented visibility","Each signal sits alone; there is no single live picture.")]
for i,(h,b) in enumerate(pains):
    col=i%2; row=i//2; bx=0.6+col*6.17; by=2.12+row*1.45
    rect(s,bx,by,5.95,1.28,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    rect(s,bx,by,5.95,0.10,fill=RED,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    text(s,bx+0.28,by+0.2,5.4,0.4,[[P(h,15.5,INK,True)]])
    text(s,bx+0.28,by+0.62,5.4,0.6,[[P(b,12,SUBTLE)]])
text(s,0.6,6.62,12.1,0.5,[[P("When an elevator is slow, costly, or insecure, the whole building feels it — today's simple controllers were not built for this.",13,INK,True,it=True)]])
notes(s,"""
How are elevators managed today? Mostly by fixed, isolated controllers - no live software
copy, no safe audited way to send remote commands, and security added as an afterthought.
That causes real pains. Maintenance is reactive: we only see a fault after it happens.
Waiting time is long because dispatch (the strategy that serves calls) is clock-based, not
context-aware. Electricity bills are high with no energy strategy. Maintenance costs are
high because breakdowns are unplanned. Security and access are weak and siloed. And
visibility is fragmented - no single live picture. These are exactly the pains our platform
is designed to address. [~50s] -> So here is our solution.
""")

# =========================================================================== #
# 4 — OUR SOLUTION (why better)
# =========================================================================== #
s=slide(); title_block(s,"Our Solution","One platform — and why it beats a simple elevator",4)
rows=[("State","local only, no memory","a live digital twin + full history"),
      ("Dispatch","fixed, clock-based","context-aware policies (less waiting)"),
      ("Energy","no strategy","ECO policy + energy analytics"),
      ("Maintenance","reactive, costly","predictive, planned work orders"),
      ("Security","limited, siloed","RFID + access rules + audited commands"),
      ("Remote control","none / unsafe","safety-gated commands, fully logged"),
      ("Visibility","scattered","one screen, one source of truth")]
tx,ty,tw=0.6,1.72,8.0
rect(s,tx,ty,tw,0.42,fill=INK)
text(s,tx+0.15,ty+0.04,2.0,0.36,[[P("Aspect",11,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
text(s,tx+2.1,ty+0.04,2.7,0.36,[[P("Simple elevator",11,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
text(s,tx+4.9,ty+0.04,3.0,0.36,[[P("Our platform",11,CYAN,True)]],anchor=MSO_ANCHOR.MIDDLE)
y=ty+0.42
for i,(a,b,c) in enumerate(rows):
    bg=CARD if i%2==0 else GREY_BG
    rect(s,tx,y,tw,0.56,fill=bg,line=RULE,lw=0.5)
    text(s,tx+0.15,y+0.03,1.95,0.5,[[P(a,11,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,tx+2.1,y+0.03,2.7,0.5,[[P(b,10.5,SUBTLE)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,tx+4.9,y+0.03,3.0,0.5,[[P(c,10.5,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.56
# right takeaways
rx=8.95
take=[("Software is the core","The platform is the contribution. It is built to drop onto a real elevator — only the device under it changes.",CYAN_DK,CYAN_BG),
      ("Safe by design","AI never controls the elevator. A fixed safety check approves every command.",RED,RED_BG),
      ("Reproducible","All open, self-hosted tools — no mandatory cloud. Anyone can rebuild it.",GREEN,GREEN_BG)]
yy=1.72
for h,b,col,bg in take:
    rect(s,rx,yy,3.78,1.55,fill=bg,line=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.07)
    text(s,rx+0.22,yy+0.14,3.4,0.4,[[P(h,14,col,True)]])
    text(s,rx+0.22,yy+0.55,3.4,0.95,[[P(b,11.5,INK)]])
    yy+=1.66
notes(s,"""
Our answer is one integrated platform, and it beats a simple elevator on every line. Today
the state is local with no memory - we give a live digital twin (a live software copy) plus
full history. Dispatch is fixed - ours is context-aware, so less waiting. We add an energy
strategy, predictive maintenance instead of costly reactive repairs, real access control
with audited commands, safe remote control, and one single live picture. Three things to
remember, on the right. The software is the core contribution - it is built to drop onto a
real elevator. It is safe by design: AI never controls the elevator, a fixed check approves
every command. And it is fully reproducible with open tools. Important honesty: we are
designed to reduce waiting, energy and maintenance costs - we give the tools to do it; we
do not claim measured savings on a real building yet. [~60s] -> What exactly were we asking?
""")

# =========================================================================== #
# 5 — RESEARCH QUESTION & HYPOTHESIS
# =========================================================================== #
s=slide(); title_block(s,"Research Question & Hypothesis","Improve safety — without handing control to AI",5)
rect(s,0.6,1.72,12.13,1.5,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
rect(s,0.6,1.72,0.12,1.5,fill=CYAN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
text(s,0.95,1.88,11.4,0.3,[[P("RESEARCH QUESTION",11.5,CYAN_DK,True)]])
text(s,0.95,2.22,11.4,1.0,[[P("Can a small elevator prototype — linked to a digital twin and a fixed-rule multi-agent layer — measurably improve ",14,INK),
    P("traceability, command safety and predictive-maintenance support",14,INK,True),
    P(", without giving AI any control over the elevator?",14,INK)]])
text(s,0.6,3.42,12,0.3,[[P("HYPOTHESIS — three pillars make the platform safe by design",12,STEEL,True)]])
tri=[("1  Authoritative state","Eclipse Ditto holds the one true, up-to-date copy (the source of truth).",CYAN_DK,CYAN_BG),
     ("2  Mandatory gate","Every command passes a fixed safety check before any change.",RED,RED_BG),
     ("3  Advisory learning","The LLM explains only; the ML brain stays in shadow.",GREEN,GREEN_BG)]
x=0.6
for h,b,col,bg in tri:
    rect(s,x,3.78,3.84,1.55,fill=bg,line=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.07)
    text(s,x+0.24,3.96,3.4,0.4,[[P(h,15,col,True)]])
    text(s,x+0.24,4.42,3.4,0.85,[[P(b,12,INK)]]); x+=4.145
rect(s,0.6,5.55,12.13,0.82,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
text(s,0.9,5.55,11.6,0.82,[[P("⇒  Traceable, auditable, safe-by-design supervision — ",13.5,WHITE,True),
    P("proven in software, on its own, apart from the physical elevator it protects.",13.5,RGBColor(0xCF,0xF6,0xFB))]],anchor=MSO_ANCHOR.MIDDLE)
notes(s,"""
Our question, in plain words: can a small elevator prototype, linked to a digital twin and
a fixed-rule multi-agent layer, measurably improve traceability, command safety and
predictive-maintenance support - without giving AI any control over the elevator? Our
answer rests on three pillars. One: Eclipse Ditto holds the one true, up-to-date copy.
Two: every command passes a fixed safety check before any change. Three: learning stays
advisory - the LLM only explains, the ML brain only practises in shadow. Together these
give traceable, auditable, safe-by-design supervision that we can prove in software on its
own. [~50s]
[Likely question - does an LLM run the elevator? No. "Agentic" here means small fixed-rule
helpers; the LLM only writes explanations and the safety check never consults it.]
-> Here is what we built.
""")

# =========================================================================== #
# 6 — CONTRIBUTIONS
# =========================================================================== #
s=slide(); title_block(s,"Contributions","One reproducible architecture — seven contributions",6)
rect(s,0.6,1.66,12.13,0.6,fill=GREEN_BG,line=GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
text(s,0.85,1.74,11.6,0.45,[[P("“Agentic AI” here = ",12.5,GREEN,True),
    P("small single-job n8n helpers that follow fixed rules.  Not autonomous.  Not generative control.",12.5,INK)]],anchor=MSO_ANCHOR.MIDDLE)
contribs=["A reduced-scale 4-floor ESP32-S3 prototype — our test bench for the platform",
          "Secure messaging baseline: per-elevator IDs, sign-in, access rules, encryption (TLS)",
          "Two-way bridge: tidies live readings up, carries approved commands back",
          "Eclipse Ditto model keeping dashboard, agents and database on one shared truth",
          "Small single-job n8n agents + one fixed command safety check",
          "Adaptive dispatch: a fixed scorer runs; a learning model only practises alongside",
          "Next.js control-room dashboard that acts through safe server APIs, never raw messages"]
pos=[(0.6,2.42),(6.77,2.42),(0.6,3.6),(6.77,3.6),(0.6,4.78),(6.77,4.78),(0.6,5.96)]
for i,(bx,by) in enumerate(pos):
    rect(s,bx,by,5.96,1.08,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    o=rect(s,bx+0.18,by+0.27,0.54,0.54,fill=CYAN_BG,shape=MSO_SHAPE.OVAL)
    tf=o.text_frame; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(i+1); r.font.bold=True; r.font.size=Pt(15); r.font.name=FONT; r.font.color.rgb=CYAN_DK
    text(s,bx+0.92,by+0.13,4.9,0.85,[[P(contribs[i],11.5,INK)]],anchor=MSO_ANCHOR.MIDDLE)
notes(s,"""
Our contribution is one disciplined, reproducible architecture - not a single algorithm.
First a definition we will reuse: "agentic AI" here means small single-job n8n helpers that
follow fixed rules - not autonomous, not generative control. The seven contributions: the
ESP32-S3 prototype as our test bench; a secure messaging baseline; a two-way bridge; the
Ditto model as one shared truth; small agents plus one fixed safety check; adaptive
dispatch with the learning model kept in shadow; and a dashboard that acts only through
safe server APIs. [~45s] -> Before the deep dive, the tools in plain words.
""")

# =========================================================================== #
# 7 — THE TECHNOLOGY IN PLAIN ENGLISH
# =========================================================================== #
s=slide(); title_block(s,"The Building Blocks","The technology we used — in plain words",7)
tech=[("ESP32-S3","controller","A tiny computer inside the elevator. It reads sensors and drives the motors.","The elevator's on-board brain."),
      ("MQTT","device messaging","A lightweight way for machines to send short messages by topic.","Like radio channels for devices."),
      ("Mosquitto","message broker","The post office that receives and delivers every message.","The sorting office."),
      ("The Bridge","translator","Tidies device messages into the twin; carries approved commands back.","A translator and clerk."),
      ("Eclipse Ditto","digital twin","The always-current software copy of the elevator — the source of truth.","The elevator's live status card."),
      ("PostgreSQL","history","The database that remembers what happened over time.","The logbook or diary."),
      ("n8n","the agents","Small automated helpers — one each for analysis, safety, alerts.","An office of specialists."),
      ("Next.js","the dashboard","The web app that shows operators the screen they watch and act on.","The control-room display."),
      ("Docker","packaging","Packs each program in its own box so it starts reliably anywhere.","Shipping containers for software."),
      ("LLM / Ollama","explainer","Writes plain-English notes about what's happening. Advisory only.","A commentator, not the driver.")]
gx=0.6; cw=4.04; ch=1.18
for i,(name,tag,body,an) in enumerate(tech):
    col=i%3; row=i//3; bx=gx+col*4.12; by=1.66+row*1.26
    rect(s,bx,by,cw,ch,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    try: place(s,TECH_ICON[i],bx+0.18,by+0.16,0.34,0.34)
    except Exception: pass
    text(s,bx+0.62,by+0.13,cw-0.7,0.3,[[P(name,12.5,INK,True),P("  "+tag,9.5,STEEL)]])
    text(s,bx+0.2,by+0.5,cw-0.35,0.5,[[P(body,9.8,SUBTLE)]])
    text(s,bx+0.2,by+0.92,cw-0.35,0.22,[[P(an,9.5,CYAN_DK,it=True)]])
rect(s,0.6,6.66,12.13,0.46,fill=CYAN_BG,line=CYAN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.2)
text(s,0.6,6.66,12.13,0.46,[[P("Together these form a digital twin: a live software copy of the real elevator.",12,CYAN_DK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
notes(s,"""
Before the deep dive, the tools in plain words. The ESP32-S3 is the small computer inside
the elevator. MQTT is how it sends short messages; Mosquitto is the post office that
delivers them. A small Bridge program tidies those messages into Eclipse Ditto - the live
software copy of the elevator, our single source of truth. A database keeps the history.
n8n runs small helpers, each with one job. Next.js is the screen operators use. Docker
packs everything so it starts reliably. And the language model only writes plain-English
notes - it never controls anything. Two words to keep: Ditto is the truth, the safety check
is the guard. [~50s] -> Now let us open the platform.
""")

# =========================================================================== #
# 8 — ARCHITECTURE WHOLE
# =========================================================================== #
s=slide(); title_block(s,"Inside the Platform","The whole map — readings go up, safe commands come down",8)
picture(s,SYNOPTIC,0.5,1.66,12.33,4.7,border=True)
text(s,0.6,6.45,12.1,0.5,[[P("We now walk this map part by part. ",11.5,INK,True),
    P("For each part we ask: what we did, which tools, how it works, and why we chose it.",11.5,SUBTLE)]])
notes(s,"""
This one map is our whole platform. Read it left to right. The elevator and its ESP32-S3
publish readings over MQTT. The bridge tidies them into Eclipse Ditto, the source of truth.
The n8n agents, the database, and the dashboard all read from Ditto. Commands travel the
other way and pass the safety check first. Three rules hold everywhere: Ditto is the truth;
commands are checked requests, not direct writes; the LLM only explains. We will now walk
this map part by part - and for every part we ask what we did, which tools, how it works,
and why we chose it. [~45s] -> Part one: the field.
""")

# =========================================================================== #
# 9..16 — DEEP-DIVE LAYERS
# =========================================================================== #
def deep(idx_active, kicker, ttl, page, what, which, how, why, note,
         shot=None, logo=None, shot_caption=None):
    s=slide(); title_block(s,kicker,ttl,page); pipeline(s,idx_active)
    if logo:
        try: place(s,logo,11.55,0.62,1.0,0.7)
        except Exception: pass
    if shot:
        wwhw(s,what,which,how,why,x=0.6,y=2.45,w=6.7,ch=1.85)
        picture(s,shot,7.55,2.45,5.2,4.0,caption=shot_caption,border=True)
    else:
        wwhw(s,what,which,how,why,x=0.6,y=2.55,w=12.13,ch=1.95)
    notes(s,note); return s

deep(0,"Inside the Platform — Field","Part 1: the elevator + ESP32-S3 (our test bench)",9,
     "A reduced-scale 4-floor elevator that produces real readings (telemetry).",
     "ESP32-S3 firmware; stepper & door motors; buttons; RFID; simulated sensor proxies.",
     "Firmware reads inputs, runs a safe state machine, and publishes telemetry over MQTT.",
     "A cheap, safe way to exercise the software with realistic signals. (A small part of the work.)",
     """
Part one, the field layer. What we did: a reduced-scale four-floor elevator that produces
real readings. Which tools: an ESP32-S3 firmware, stepper and door motors, buttons, an RFID
reader, and a few simulated sensor proxies. How it works: the firmware reads the inputs,
runs a safe state machine, and publishes the readings over MQTT. Why this choice: it is a
cheap and safe way to feed the software realistic signals - and we stress, this hardware is
a small part of the work. The value is the software above it. [~45s] -> Part two: messaging.
""",
     shot=F("prototype/panels/hardware_evidence_panel.png"),
     shot_caption="Our test bench — shaft, drive, electronics, safety panel")

deep(1,"Inside the Platform — Messaging","Part 2: MQTT (Mosquitto) + the bridge",10,
     "Move readings up; carry approved commands down.",
     "Eclipse Mosquitto (broker) + a Node.js bridge (translator).",
     "Publish/subscribe topics; the bridge tidies messages into twin fields and forwards only approved commands.",
     "Lightweight and decoupled — and the security boundary (sign-in + access rules + TLS) lives here.",
     """
Part two, messaging. What we did: move readings up and carry approved commands down. Which
tools: the Mosquitto broker - the post office - and a Node.js bridge, our translator. How
it works: devices publish to topics, others subscribe; the bridge tidies each message into
the right twin field, and on the way back it forwards only approved commands. Why: MQTT is
lightweight and decoupled, and this is exactly where our security boundary sits - sign-in,
access rules, and encryption. [~45s] -> Part three: the digital twin.
""", logo=BRAND["mosq"])

deep(2,"Inside the Platform — Digital Twin","Part 3: Eclipse Ditto, the source of truth",11,
     "Keep one always-current digital copy of the elevator.",
     "Eclipse Ditto — a Thing with features (status fields).",
     "The bridge writes per-field updates; the dashboard and agents read the twin, never raw MQTT.",
     "One shared truth keeps everything consistent — and ready for many elevators.",
     """
Part three, the digital twin. What we did: keep one always-current digital copy of the
elevator. Which tool: Eclipse Ditto, which models the elevator as a Thing with feature
fields. How it works: the bridge writes small per-field updates, and the dashboard and
agents always read the twin - never the raw messages. Why: one shared truth keeps the whole
system consistent, and it is ready to grow from one cabin to many. On the right is the live
twin during our final run - fifteen feature fields, in sync. [~45s] -> Part four: the agents.
""", shot=F("screenshots/edited/ditto_explorer_thing_overview.png"),
     shot_caption="Live Ditto twin — 15 feature fields, in sync", logo=BRAND["ditto"])

deep(3,"Inside the Platform — Agentic AI","Part 4: small single-job n8n agents",12,
     "Analyse, secure, maintain, notify, audit — each a small helper.",
     "n8n workflows; fixed rule engines; an optional local LLM for plain-English notes.",
     "Agents read the twin, run fixed rules, route validated actions; one shared message format.",
     "Separation of concerns = inspectable, testable, safe — no opaque autonomy.",
     """
Part four, the agentic layer. What we did: split the smart work into small helpers - one
for analysis, one for security, one for maintenance, notifications, and audit. Which tools:
n8n workflows with fixed rule engines, plus an optional local LLM that only writes
plain-English notes. How it works: each agent reads the twin, runs fixed rules, and routes
validated actions, all in one shared message format. Why: separating the jobs makes
everything inspectable, testable and safe - there is no opaque autonomy. Even a security
lockdown goes through the same safety check as an operator command. [~50s] -> Part five:
dispatch.
""", shot=F("screenshots/edited/n8n_ingestion_workflow.png"),
     shot_caption="A real n8n agent workflow", logo=BRAND["n8n"])

deep(4,"Inside the Platform — Dispatch","Part 5: adaptive dispatch (champion / challenger)",13,
     "Pick the best serving strategy (dispatch policy) for the live situation.",
     "Brain A — a fixed scorer (active); Brain B — a learning model (shadow).",
     "Build a context vector → score the 8 policies → the active brain binds → command via the safety check.",
     "Adaptivity without risk — the learner only practises until it provably wins. (Detail next slide.)",
     """
Part five, dispatch - how the cabin decides to serve calls. What we did: pick the best
strategy for the live situation. Which tools: two brains. Brain A is a fixed scorer and is
active; Brain B is a learning model and runs only in shadow. How it works: we build a
context vector, score the eight strategies, the active brain binds the choice, and it still
goes through the safety check. Why: we get adaptivity without risk - the learner only
practises until it can prove it is better. We give the model its own slide in a moment.
[~45s] -> Part six: the safety check.
""")

# Part 6 — Command safety (full funnel emphasis)
s=slide(); title_block(s,"Inside the Platform — Safety","Part 6: the fixed command safety check",14); pipeline(s,5)
srcs=["Operator (dashboard)","n8n agent","System task"]; y=2.45
for sname in srcs:
    rect(s,0.6,y,2.55,0.55,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    text(s,0.6,y,2.55,0.55,[[P(sname,11,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    arrow(s,3.2,y+0.11,0.5,0.32,fill=STEEL); y+=0.75
text(s,0.6,4.7,2.55,0.4,[[P("all POST →",11,STEEL,True)]],align=PP_ALIGN.CENTER)
rect(s,3.85,2.4,3.5,2.75,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
text(s,4.0,2.56,3.2,0.55,[[P("DETERMINISTIC",12.5,CYAN,True)],[P("SAFETY CHECK (fixed rules)",12.5,CYAN,True)]],align=PP_ALIGN.CENTER)
yy=3.35
for pr in ["allow-listed","valid + authorised source","twin is up to date","no emergency / no overload","safe door state","risk within limit","Ditto available"]:
    text(s,4.15,yy,3.0,0.25,[[P("✓ "+pr,10.5,WHITE)]]); yy+=0.255
rect(s,7.75,2.45,4.95,1.25,fill=GREEN_BG,line=GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,7.95,2.57,4.6,0.3,[[P("ACCEPTED",13,GREEN,True)]])
text(s,7.95,2.92,4.6,0.7,[[P("update the twin → bridge → device → built-in safety stop",12,INK)]])
arrow(s,7.4,2.9,0.4,0.32,fill=GREEN)
rect(s,7.75,3.9,4.95,1.25,fill=RED_BG,line=RED,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,7.95,4.02,4.6,0.3,[[P("REJECTED",13,RED,True)]])
text(s,7.95,4.37,4.6,0.7,[[P("logged with a reason → nothing sent, no twin change",12,INK,True)]])
arrow(s,7.4,4.35,0.4,0.32,fill=RED)
rect(s,0.6,5.45,12.13,0.78,fill=RED,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
text(s,0.6,5.45,12.13,0.78,[[P("WHY it matters: rejected commands change nothing — verified 33/33 in tests + 9 live decisions. This is how we keep AI from authorising unsafe actions.",13,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
text(s,0.6,6.28,12.13,0.4,[[P("The LLM is completely outside this decision — the check never asks it and never takes it as permission.",11,SUBTLE,it=True)]],align=PP_ALIGN.CENTER)
notes(s,"""
Part six is the heart of the work: the fixed command safety check. What we did: approve or
block every command with fixed rules. Every command - from an operator, an agent, or a
system task - posts into one check. It must be allow-listed, valid, from an authorised
source, the twin must be fresh, no emergency or overload, a safe door state, risk within
limit, and Ditto available. If all pass, we update the twin, the bridge forwards it, and
the firmware still applies its own safety stops. If anything fails, we reject with a reason
and change nothing - zero twin writes. We verified that 33 out of 33 in tests and on 9 live
decisions. And the LLM is completely outside this decision. This is how an AI-assisted
system is kept from authorising unsafe actions. [~60s]
[Likely question - how do you stop the AI doing something unsafe? This check. The AI may
propose; only the fixed rules decide; the firmware is the last line.] -> Part seven: memory.
""")

deep(6,"Inside the Platform — Data","Part 7: PostgreSQL / TimescaleDB (memory)",15,
     "Remember everything over time.",
     "PostgreSQL + the TimescaleDB extension (time-series).",
     "Stores telemetry, audit, command & dispatch decisions, work orders, notifications.",
     "Traceability — and the data foundation for analytics and future learning.",
     """
Part seven, the memory. What we did: remember everything over time. Which tool: PostgreSQL
with the TimescaleDB extension for time-series data. How it works: it stores the telemetry
history, the audit trail, every command and dispatch decision, the maintenance work orders,
and the notifications. Why: this gives us full traceability - we can always show what
happened and why - and it is the data foundation for analytics and for training the
learning model later. During our run it held over nineteen thousand six hundred telemetry
rows. [~40s] -> Part eight: the screen.
""", logo=BRAND["pg"])

deep(7,"Inside the Platform — Screen","Part 8: the Next.js SCADA dashboard",16,
     "Give operators one control-room screen.",
     "Next.js web app; server-side history APIs.",
     "Reads the twin (live updates); sends actions only through safe server APIs.",
     "Supervision and command requests — never direct device control from the browser.",
     """
Part eight, the screen. What we did: give operators one control-room screen. Which tool: a
Next.js web app with server-side history APIs. How it works: it reads the live twin and
shows the full picture, and it sends actions only through safe server APIs - the browser
itself can never command the device. Why: the dashboard is for supervision and for sending
command requests, not for direct control. That keeps the safety boundary intact. On the
right is our Digital Twin page showing live state from Ditto. [~40s] -> Now the
machine-learning model itself.
""", shot=F("screenshots/edited/dashboard_digital_twin_page.png"),
     shot_caption="ElevatorOS — live state from the twin", logo=BRAND["next"])

# =========================================================================== #
# 17 — THE ML MODEL
# =========================================================================== #
s=slide(); title_block(s,"The Machine-Learning Model","The learning model we are building — honestly",17)
# left column
lx=0.6; lw=6.5
rect(s,lx,1.66,lw,5.05,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,lx+0.25,1.8,lw-0.5,0.3,[[P("WHAT IT IS",11.5,CYAN_DK,True)]])
text(s,lx+0.25,2.12,lw-0.5,0.85,[[P("Brain B = ",12,INK),P("ml_v1",12,CYAN_DK,True,f=MONO),
    P(", a softmax linear model (multinomial logistic regression). From the live context it gives each of the 8 dispatch strategies a probability and picks the best.",12,INK)]])
text(s,lx+0.25,3.05,lw-0.5,0.5,[[P("It is the ",11.5,INK),P("challenger",11.5,INK,True),
    P("; Brain A (a fixed scorer) is the ",11.5,INK),P("champion",11.5,INK,True),P(" and stays active.",11.5,INK)]])
text(s,lx+0.25,3.62,lw-0.5,0.3,[[P("INPUTS (the shared feature vector)",11.5,STEEL,True)]])
for i,t in enumerate(["Traffic: up/down-call ratio, lobby trend, queue, starvation",
                      "Energy: tariff flag, power-vs-baseline, kWh budget",
                      "Machine health: motor temp, vibration, RUL (life left)",
                      "Load + security state; plus a fairness rule"]):
    text(s,lx+0.25,3.95+i*0.32,lw-0.5,0.3,[[P("•  ",11,CYAN_DK,True),P(t,11,SUBTLE)]])
text(s,lx+0.25,5.32,lw-0.5,0.3,[[P("HOW IT LEARNS & IS PROMOTED",11.5,STEEL,True)]])
text(s,lx+0.25,5.64,lw-0.5,1.0,[[P("Trained offline (~4,000 samples), evaluated (~2,000). A shared reward scores outcomes. It is switched on ",11,INK),
    P("only if",11,INK,True),P(" it provably beats Brain A ",11,INK),P("and a human approves",11,INK,True),
    P(". One flag switches; rollback is instant.",11,INK)]])
# right column
rx=7.35; rw=5.38
rect(s,rx,1.66,rw,2.0,fill=GREEN_BG,line=GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
text(s,rx+0.25,1.78,rw-0.5,0.3,[[P("PROPERTIES",11.5,GREEN,True)]])
for i,t in enumerate(["Explainable — per-feature weight attributions","Safe — overrides pre-empt it; the safety check still applies","Auditable — every decision + outcome is stored","Shadow mode — it never acts on the elevator"]):
    text(s,rx+0.25,2.12+i*0.36,rw-0.5,0.32,[[P("•  ",11,GREEN,True),P(t,11,INK)]])
rect(s,rx,3.8,rw,2.9,fill=AMBER_BG,line=AMBER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
text(s,rx+0.25,3.92,rw-0.5,0.3,[[P("HONEST RESULTS",11.5,AMBER,True)]])
text(s,rx+0.25,4.26,rw-0.5,0.85,[[P("Current: ",11.5,AMBER,True),
    P("Brain B imitates Brain A at about 66%. Its reward does not clear the gate — so we correctly keep Brain A active. Brain B stays shadow-only.",11.5,INK)]])
text(s,rx+0.25,5.18,rw-0.5,0.55,[[P("Value: ",11.5,GREEN,True),
    P("it proves the full pipeline — train, evaluate, gate, promote, roll back — works safely.",11.5,INK)]])
text(s,rx+0.25,5.86,rw-0.5,0.8,[[P("Future: ",11.5,CYAN_DK,True),
    P("with real outcome data + a richer model, a challenger that truly cuts waiting and energy could be promoted — only after it wins and a human approves.",11.5,INK)]])
notes(s,"""
Now the learning model itself, told honestly. What it is: Brain B, called ml_v1, is a
softmax linear model - a multinomial logistic regression. From the live situation it gives
each of the eight dispatch strategies a probability and picks the best. It is the
challenger; Brain A, the fixed scorer, is the champion and stays active. Its inputs are a
shared feature vector: traffic, energy, machine health including remaining useful life,
load and security, plus a fairness rule. How it learns: trained and evaluated offline on a
shared reward; it is switched on only if it provably beats Brain A and a human approves -
one flag, instant rollback - and it always runs in shadow, never acting on the elevator.
Honest results: today it imitates Brain A at about sixty-six percent; a linear model cannot
capture every rule, its reward does not clear the gate, so we correctly keep Brain A
active. Its value is that it proves the whole pipeline works safely. In future, with real
outcome data and a richer model, a challenger that truly cuts waiting and energy could be
promoted - only after it wins and a human approves. [~65s] -> What did we actually validate?
""")

# =========================================================================== #
# 18 — VALIDATION MATRIX
# =========================================================================== #
s=slide(); title_block(s,"Validation & Evidence","Software proven; the hardware boundary stated plainly",18)
text(s,0.6,1.6,12.1,0.4,[[P("A test passes only when real evidence exists (logs, screenshots, queries, photos).  ",12,INK,True),
    P("Software proof is not hardware proof.",12,RED,True)]])
data=[("Command safety check","33/33 assertions + 9 live decisions","PASS",GREEN,GREEN_BG),
      ("Adaptive dispatch","61 assertions + live policy adoption","PASS",GREEN,GREEN_BG),
      ("Bridge ↔ Ditto sync","Bridge logs + live Ditto exports","PASS",GREEN,GREEN_BG),
      ("Messaging security","Allowed/denied access + sign-in captures","PASS",GREEN,GREEN_BG),
      ("Database / history","7 tables, >19,600 telemetry rows","PASS",GREEN,GREEN_BG),
      ("Dashboard","type-check clean + page screenshots","PASS",GREEN,GREEN_BG),
      ("n8n workflows","validator + published-workflow figures","PASS",GREEN,GREEN_BG),
      ("ESP32 firmware path","broker-side captures (no serial log)","PARTIAL",AMBER,AMBER_BG),
      ("Hardware bench","photos, dimensions, ~6 s/floor","DOCUMENTED",AMBER,AMBER_BG),
      ("KY-024 / SPDT / RFID","mounted / firmware only","OUT OF SCOPE",GREY,GREY_BG)]
for i,(name,ev,st,col,bg) in enumerate(data):
    col_i=i//5; row=i%5; bx=0.6+col_i*6.18; by=2.15+row*0.9
    rect(s,bx,by,5.95,0.78,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    rect(s,bx,by,0.09,0.78,fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    text(s,bx+0.25,by+0.08,3.6,0.32,[[P(name,11.5,INK,True)]])
    text(s,bx+0.25,by+0.4,4.4,0.32,[[P(ev,10,SUBTLE)]])
    chip(s,bx+4.5,by+0.24,st,bg,col,size=9)
notes(s,"""
What did we actually validate? Our rule is on the slide: a test passes only when real
evidence exists - logs, screenshots, queries, photos - and software proof is not hardware
proof. In green, all software-evidenced on our live run: the safety check, adaptive
dispatch, bridge-to-Ditto sync, messaging security, the database with over nineteen
thousand six hundred rows, the dashboard, and the n8n workflows. In amber: the ESP32 path
is broker-side evidenced without a serial log, and the hardware bench is documented by
photos and dimensions. In grey, openly out of scope: the KY-024 floor sensors, the SPDT
door switches, and physical RFID. We never dress a grey item up as green. [~55s] -> The
headline numbers.
""")

# =========================================================================== #
# 19 — NUMBERS + EVIDENCE
# =========================================================================== #
s=slide(); title_block(s,"Validation & Evidence","The numbers — and the live evidence behind them",19)
tiles=[("127","automated tests\n(94 Node + 33 Python)"),("33/33","safety-check assertions\n+ 9 live decisions"),
       (">19,600","telemetry rows\nsince 2026-05-20"),("15","live twin fields\n(twin in sync)"),
       ("2026-06-13","final live run on\nthe full Docker stack")]
x=0.6
for big,small in tiles:
    rect(s,x,1.72,2.3,1.5,fill=CARD,line=RULE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.07)
    rect(s,x,1.72,2.3,0.1,fill=CYAN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    text(s,x+0.1,1.92,2.1,0.6,[[P(big,25,CYAN_DK,True)]],align=PP_ALIGN.CENTER)
    ls=small.split("\n")
    text(s,x+0.1,2.6,2.1,0.6,[[P(ls[0],10.5,SUBTLE)],[P(ls[1],10.5,SUBTLE)]],align=PP_ALIGN.CENTER)
    x+=2.43
picture(s,F("screenshots/edited/dashboard_digital_twin_page.png"),0.6,3.55,6.0,3.05,
        caption="ElevatorOS Digital Twin — live state from the twin")
picture(s,F("screenshots/panels/dashboard_ai_security_maintenance_panel.png"),6.9,3.55,5.83,3.05,
        caption="AI · Security · Maintenance (RFID = twin state, not a bench test)")
notes(s,"""
The numbers, all from our validation. 127 automated tests pass - 94 Node and 33 Python. The
safety check is 33 out of 33 plus nine live decisions. Over nineteen thousand six hundred
telemetry rows since the twentieth of May. Fifteen live twin fields, in sync. And one final
live run on the full Docker stack on the thirteenth of June. Below are two pieces of that
evidence - the live Digital Twin page, and the combined AI, security and maintenance view.
We label the second carefully: the RFID state there is read from the twin, it is not a
physical reader bench test. [~50s] -> Now the limits and what comes next.
""")

# =========================================================================== #
# 20 — LIMITATIONS
# =========================================================================== #
s=slide(); title_block(s,"Limitations & Future Work","Every limit, paired with its hardened next step",20)
rows=[("KY-024 floor confirmation","Alignment, debouncing, homing, optional encoder feedback"),
      ("SPDT door confirmation","Switch wiring, detection logic, obstruction sensing"),
      ("Simulated load (turn-knob)","Load cell + HX711 amplifier (calibrated)"),
      ("Simulated temp / vibration","Real sensors with calibration"),
      ("RFID with no bench test yet","MFRC522 UID authorisation + captured scans"),
      ("Learning model only practises","Promote only after real-outcome evidence + human review"),
      ("Prototype safety logic","Certified safety hardware + formal hazard analysis")]
tx,ty,tw=0.6,1.9,12.13
rect(s,tx,ty,tw,0.44,fill=INK)
text(s,tx+0.2,ty+0.04,4.5,0.36,[[P("Current (this work)",12,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
text(s,tx+6.15,ty+0.04,6,0.36,[[P("Future hardened next step",12,CYAN,True)]],anchor=MSO_ANCHOR.MIDDLE)
y=ty+0.44
for j,(a,b) in enumerate(rows):
    bg=CARD if j%2==0 else GREY_BG
    rect(s,tx,y,tw,0.55,fill=bg,line=RULE,lw=0.5)
    chip(s,tx+0.15,y+0.13,"now",GREY_BG,GREY,size=9)
    text(s,tx+0.85,y+0.03,4.95,0.5,[[P(a,11.5,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    arrow(s,tx+5.75,y+0.17,0.35,0.22,fill=CYAN)
    text(s,tx+6.25,y+0.03,5.7,0.5,[[P(b,11.5,SUBTLE)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.55
text(s,tx,y+0.08,tw,0.4,[[P("Two stages: ",12,STEEL,True),
    P("soon — real hardware evidence (videos, logs, sensors);   later — production hardening and the bigger picture.",11.5,SUBTLE)]])
notes(s,"""
We frame limits as a roadmap, each paired with its fix. The KY-024 floor sensors and SPDT
door switches need alignment and detection logic. The simulated load, temperature and
vibration need real calibrated sensors. RFID needs a real reader with captured scans. The
learning model should be promoted only after real outcome evidence and human review. And
the prototype safety logic would need certified safety hardware before any real building.
Two stages: soon, the physical hardware evidence; later, production hardening and the bigger
picture - which is our next slide. [~45s] -> Where this can go.
""")

# =========================================================================== #
# 21 — INTEGRATION / BIGGER PICTURE
# =========================================================================== #
s=slide(); title_block(s,"The Bigger Picture","Designed to plug into building & city systems",21)
steps=[("One elevator","Live twin + history for a single cabin",CYAN_BG,CYAN_DK),
       ("Building (BMS)","Feeds the building's central monitoring — energy, security, maintenance",GREEN_BG,GREEN),
       ("Smart city","Many buildings, many twins — one dashboard, one pattern",AMBER_BG,AMBER)]
x=0.6
for i,(h,b,bg,col) in enumerate(steps):
    rect(s,x,2.2,3.7,2.4,fill=bg,line=col,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    text(s,x+0.3,2.45,3.1,0.5,[[P(h,17,col,True)]])
    text(s,x+0.3,3.05,3.1,1.4,[[P(b,13,INK)]])
    if i<2: arrow(s,x+3.78,3.2,0.5,0.42,fill=STEEL)
    x+=4.15
text(s,0.6,5.1,12.1,0.4,[[P("Already built in: ",12.5,INK,True),
    P("per-elevator IDs, a standard twin model, and open web APIs — so the same platform scales out.",12.5,SUBTLE)]])
rect(s,0.6,5.7,12.13,0.7,fill=AMBER_BG,line=AMBER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
text(s,0.85,5.7,11.7,0.7,[[P("Honesty: ",12,AMBER,True),
    P("this is an architectural capability and roadmap, shown for one cabin. Multi-building deployment is future work.",12,INK)]],anchor=MSO_ANCHOR.MIDDLE)
notes(s,"""
Where can this go? Because we already use per-elevator IDs, a standard twin model, and open
web APIs, the same platform is designed to plug into bigger systems. One elevator becomes
one managed asset in a Building Management System - feeding the building's central energy,
security and maintenance monitoring. And the same pattern scales to a smart city: many
buildings, many twins, one dashboard. To stay honest: this is an architectural capability
and a roadmap. We have shown it for one cabin; multi-building deployment is future work.
[~40s] -> To conclude.
""")

# =========================================================================== #
# 22 — CONCLUSION
# =========================================================================== #
s=slide(); title_block(s,"Conclusion","The research question — answered",22)
rect(s,0.6,1.72,12.13,1.45,fill=CYAN_BG,line=CYAN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
text(s,0.95,1.9,11.5,1.2,[[P("Yes — ",16,CYAN_DK,True),
    P("traceability, command safety and predictive-maintenance support are all improved, shown in software, ",15,INK),
    P("without giving AI any control of the elevator",15,INK,True),
    P(" — while physical-hardware validation remains a stated boundary.",15,INK)]])
text(s,0.6,3.35,12.1,0.35,[[P("The contribution is the whole safe-by-design software architecture — not any single algorithm.",13,INK,True)]])
pil=[("Traceability","Ditto as the one truth + a full audit of every decision",CYAN_DK,CYAN_BG),
     ("Command safety","A fixed safety check; rejected commands change nothing",RED,RED_BG),
     ("Maintenance support","Checkable fixed-rule risk and wear — nothing overclaimed",GREEN,GREEN_BG)]
x=0.6
for h,b,col,bg in pil:
    rect(s,x,3.85,3.94,1.5,fill=bg,line=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.07)
    text(s,x+0.25,4.02,3.5,0.4,[[P(h,15,col,True)]])
    text(s,x+0.25,4.45,3.5,0.85,[[P(b,12,INK)]]); x+=4.105
rect(s,0.6,5.6,12.13,0.72,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
text(s,0.6,5.6,12.13,0.72,[[P("A software platform proven in the lab — an honest, solid base to grow toward a hardened industrial system.",13.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
notes(s,"""
To conclude, the answer to our question is yes - with care. Traceability, command safety
and predictive-maintenance support are all improved, shown in software, without giving AI
any control of the elevator - while physical-hardware validation stays a clearly stated
boundary. Our contribution is the whole safe-by-design software architecture, not any
single algorithm: traceability from Ditto plus a full audit; command safety from the fixed
check where rejected commands change nothing; and maintenance support from checkable
fixed-rule analysis that we do not overclaim. It is a software platform proven in the lab,
and an honest, solid base to grow toward a hardened industrial system. Thank you - we
welcome your questions. [~50s]
""")

# =========================================================================== #
# 23 — THANK YOU
# =========================================================================== #
s=slide(); rect(s,0,0,0.22,7.5,fill=CYAN)
place(s,UNIV,0.7,0.6,1.15,1.1)
text(s,0.8,2.7,11,1.1,[[P("Thank you",44,INK,True)]])
text(s,0.82,3.8,11,0.6,[[P("Questions & discussion",22,CYAN_DK,True)]])
rect(s,0.85,4.6,2.0,0.05,fill=CYAN)
text(s,0.8,4.85,11.5,1.0,[[P("Agentic AI-Driven Digital Twin for Smart and Secure Elevator Management System",13.5,SUBTLE,True)],
                          [P(AUTHORS+"   •   Supervisor: "+SUP+"   •   2026",12.5,STEEL)]])
text(s,0.8,6.2,11.5,0.5,[[P("Backup slides: MQTT payloads · electrical schematic · deployment.",11,STEEL,it=True)]])
notes(s,"Thank you. We are happy to take questions. We have backup slides on the MQTT payloads, the electrical schematic, and the deployment if useful.")

# =========================================================================== #
# 24 — APPENDIX: MQTT
# =========================================================================== #
s=slide(); title_block(s,"Appendix","MQTT topics & representative payloads",23)
text(s,0.6,1.7,12,0.3,[[P("Canonical topics  —  elevator/{mqtt_safe_thing_id}/{telemetry | events | commands | status}",12.5,INK,True,f=MONO)]])
tel='''{
  "thingId": "building:floor1:elevator",
  "value": {
    "cabin": { "properties": { "current_floor": 1, "target_floor": 3,
                "direction": "UP", "load_kg": 120.0, "emergency_stop": false } },
    "door": { "properties": { "state": "CLOSED", "cycle_count": 7 } },
    "motor": { "properties": { "temperature_c": 42.0, "health_status": "GOOD" } },
    "fan": { "properties": { "state": "ON", "mode": "AUTO" } }
  }
}'''
cmd='''{
  "command": "MOVE_TO_FLOOR",
  "command_id": "CMD-2026-0001",
  "correlation_id": "CID-2026-0001",
  "thing_id": "building:floor1:elevator",
  "target_floor": 3,
  "source": "ditto_command_intent"
}'''
def codebox(x,y,w,h,title,code):
    rect(s,x,y,w,h,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.03)
    text(s,x+0.2,y+0.12,w-0.4,0.3,[[P(title,11,CYAN,True)]])
    paras=[[P(line if line else " ",9.5,RGBColor(0xE2,0xE8,0xF0),f=MONO)] for line in code.split("\n")]
    text(s,x+0.2,y+0.5,w-0.4,h-0.6,paras,line_spacing=1.0,space_after=0)
codebox(0.6,2.15,7.4,4.4,"telemetry (ESP32-S3 → MQTT → bridge → Ditto)",tel)
codebox(8.2,2.15,4.53,4.4,"approved command (Ditto intent → device)",cmd)
notes(s,"Backup: the canonical topic naming and representative telemetry and command payloads, aligned to the Ditto feature tree. Values are structural examples, not measurements.")

# =========================================================================== #
# 25 — APPENDIX: SCHEMATIC
# =========================================================================== #
s=slide(); title_block(s,"Appendix","General electrical schematic (laboratory prototype)",24)
picture(s,F("electrical/elevator_general_electrical_schematic_A4_raster.png"),0.6,1.7,12.13,4.8,
        caption="ESP32-S3 central · ATX + KCD1 PS-ON · separated rails, common ground · not certified industrial wiring")
notes(s,"Backup: our electrical schematic - ESP32-S3 central, separated logic and actuator rails on a common ground, ATX supply with a KCD1 PS-ON switch. It documents the lab design; it is not certified industrial wiring.")

# =========================================================================== #
# 26 — APPENDIX: DEPLOYMENT
# =========================================================================== #
s=slide(); title_block(s,"Appendix","Deployment topology (Docker Compose + external Ditto)",25)
svcs=[("Mosquitto","elevator-mqtt","Broker: sign-in, access rules, TLS, healthcheck"),
      ("Bridge","elevator_bridge","MQTT↔Ditto sync + command forwarding"),
      ("n8n","elevator_agents","Workflow automation / agent orchestration"),
      ("PostgreSQL/TimescaleDB","elevator_db","Telemetry, audit, commands, work orders"),
      ("Simulator","elevator_simulator","Opt-in validation & fault injection"),
      ("Ollama","elevator_ollama","Optional local explanation model (ai profile)")]
tx,ty,tw=0.6,1.85,12.13
rect(s,tx,ty,tw,0.42,fill=INK)
for cx,lab,w in [(tx,"Service",3.0),(tx+3.0,"Container",3.3),(tx+6.3,"Role",5.83)]:
    text(s,cx+0.15,ty+0.03,w-0.2,0.36,[[P(lab,12,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
y=ty+0.42
for a,b,c in svcs:
    rect(s,tx,y,tw,0.6,fill=CARD,line=RULE,lw=0.5)
    text(s,tx+0.15,y,2.85,0.6,[[P(a,11.5,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,tx+3.15,y,3.1,0.6,[[P(b,11,CYAN_DK,f=MONO)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,tx+6.45,y,5.6,0.6,[[P(c,11,SUBTLE)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.6
text(s,tx,y+0.12,tw,0.6,[[P("Runs outside Compose: ",11.5,STEEL,True),
    P("Eclipse Ditto, the Next.js dashboard, and the dispatch engine.  All free & self-hostable — no mandatory cloud.",11.5,SUBTLE)]])
notes(s,"Backup: our local deployment. Mosquitto, the bridge, n8n and TimescaleDB run under Docker Compose; Ditto, the dashboard and the dispatch engine run alongside. Everything is free and self-hostable.")

# --------------------------------------------------------------------------- #
prs.save(OUT)
print("Saved:",OUT); print("Slides:",len(prs.slides._sldIdLst)); print("Font:",FONT)
