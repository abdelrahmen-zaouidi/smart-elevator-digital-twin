#!/usr/bin/env python3
"""Text-only speaker-notes PDF (no slide images). Each slide = a heading + its notes,
flowing across A4 pages. Reads notes from the current .pptx. Run:
    python scripts/build_notes_text_pdf.py"""
import os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

ROOT = r"c:/Users/Administrator/smart-elevator-twin"
PPTX = os.path.join(ROOT, "Agentic_Elevator_Digital_Twin_Defense.pptx")
OUT  = os.path.join(ROOT, "Agentic_Elevator_Digital_Twin_Defense_NOTES_textonly.pdf")

W, H = 1240, 1754; M = 90
INK=(15,23,42); BODY=(31,41,55); SUB=(100,116,139); CYAN=(14,136,153)
RULE=(226,232,240); GREY=(148,163,184)
FB=r"C:/Windows/Fonts/segoeuib.ttf"; FR=r"C:/Windows/Fonts/segoeui.ttf"; FI=r"C:/Windows/Fonts/segoeuii.ttf"
def f(p,s): return ImageFont.truetype(p,s)

prs=Presentation(PPTX); slides=list(prs.slides)

def title(s):
    t=[]
    for sh in s.shapes:
        if sh.has_text_frame:
            x=sh.text_frame.text.strip()
            if x and "Agentic AI-Driven Elevator Digital Twin" not in x and not x.isdigit():
                t.append(x.replace("\n"," "))
    if not t: return ""
    if len(t)>1 and len(t[0])<40 and t[0].upper()==t[0]:
        return f"{t[0]} — {t[1]}"
    return t[0]

# build a measuring context
_meas=ImageDraw.Draw(Image.new("RGB",(10,10)))
def wrap(text,fnt,maxw):
    out=[]
    for para in text.split("\n"):
        para=para.strip()
        if not para: out.append(""); continue
        line=""
        for w in para.split(" "):
            tr=(line+" "+w).strip()
            if _meas.textlength(tr,font=fnt)<=maxw: line=tr
            else:
                if line: out.append(line)
                line=w
        if line: out.append(line)
    return out

pages=[]; page=None; d=None; y=0
def new_page():
    global page,d,y
    page=Image.new("RGB",(W,H),"white"); d=ImageDraw.Draw(page)
    d.rectangle([0,0,W,8],fill=CYAN)
    d.text((M,34),"Defence speaker notes  —  Zaouidi & Bengherbia",font=f(FB,18),fill=SUB)
    d.text((W-M-90,34),f"{len(slides)} slides",font=f(FR,18),fill=GREY)
    pages.append(page); y=84
new_page()

HEAD=f(FB,23); BODY_F=f(FR,23); ITAL=f(FI,22); LBL=f(FB,16)
lh=int(23*1.34)
for i,s in enumerate(slides,1):
    raw=s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
    head=f"Slide {i}"+(f" — {title(s)}" if title(s) else "")
    head_lines=wrap(head,HEAD,W-2*M)
    body_lines=wrap(raw,BODY_F,W-2*M) if raw else ["(no speaker notes)"]
    block_h=14+len(head_lines)*30+8+len(body_lines)*lh+26
    # keep heading with at least 2 body lines together; else new page
    if y+ (14+len(head_lines)*30+8+2*lh) > H-M:
        new_page()
    d.line([M,y,W-M,y],fill=RULE,width=2); y+=14
    for ln in head_lines:
        d.text((M,y),ln,font=HEAD,fill=CYAN); y+=30
    y+=8
    bf = ITAL if not raw else BODY_F
    bc = GREY if not raw else BODY
    for ln in body_lines:
        if y+lh>H-M: new_page()
        d.text((M,y),ln,font=bf,fill=bc); y+=lh
    y+=26

pages[0].save(OUT,"PDF",save_all=True,append_images=pages[1:],resolution=150.0)
print("Saved:",OUT,"| pages:",len(pages))
